"""文档解析器：支持 PDF/PPTX/DOCX/TXT，返回纯文本。"""
import os


def extract_text(path: str) -> str:
    """从文件中提取纯文本内容。支持 PDF/PPTX/DOCX/TXT。"""
    if not os.path.exists(path):
        raise FileNotFoundError(f"文件不存在: {path}")
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(path)
    elif ext in (".docx", ".doc"):
        return _extract_docx(path)
    elif ext == ".txt":
        with open(path, encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        raise ValueError(f"不支持的文件类型: {ext}（支持 PDF/PPTX/DOCX/TXT）")


def _extract_pdf(path: str) -> str:
    import pypdf
    reader = pypdf.PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    return "\n".join(pages)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            texts.append(f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def _extract_docx(path: str) -> str:
    import docx
    doc = docx.Document(path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(paragraphs)
