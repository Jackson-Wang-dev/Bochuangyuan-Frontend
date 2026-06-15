"""
Text extraction from uploaded files.

Priority:
  1. MinerU CLI (if installed) — best quality for PDF/image-heavy docs.
  2. python-docx  (.docx)
  3. python-pptx  (.pptx)
  4. pypdf        (.pdf)  — fallback
  5. Plain decode  (.txt / .md)
  6. Vision model via instructor  (.jpg / .jpeg / .png / .webp)
"""
from __future__ import annotations

import asyncio
import base64
import os
import shutil
import tempfile
from pathlib import Path


async def extract_text(file_bytes: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()

    # ── MinerU (best quality, optional) ─────────────────────────────────────
    if shutil.which("mineru"):
        return await _mineru(file_bytes, suffix)

    # ── Format-specific fallbacks (sync libs run in a thread) ────────────────
    loop = asyncio.get_event_loop()
    if suffix == ".docx":
        return await loop.run_in_executor(None, _docx, file_bytes)
    if suffix == ".pptx":
        return await loop.run_in_executor(None, _pptx, file_bytes)
    if suffix == ".pdf":
        return await loop.run_in_executor(None, _pdf, file_bytes)
    if suffix in (".txt", ".md"):
        return file_bytes.decode("utf-8", errors="replace")
    if suffix in (".jpg", ".jpeg", ".png", ".webp"):
        return await _image_vision(file_bytes, suffix)

    raise ValueError(f"不支持的文件格式：{suffix}；支持 pdf / docx / pptx / txt / jpg / png")


# --------------------------------------------------------------------------- #
# MinerU
# --------------------------------------------------------------------------- #

async def _mineru(file_bytes: bytes, suffix: str) -> str:
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / f"input{suffix}"
        out = Path(tmp) / "out"
        src.write_bytes(file_bytes)
        out.mkdir()

        # -b pipeline: CPU-only mode (no GPU required)
        cmd = ["mineru", "-p", str(src), "-o", str(out), "-b", "pipeline"]
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        _, stderr = await asyncio.wait_for(proc.communicate(), timeout=300)

        if proc.returncode != 0:
            raise RuntimeError(f"MinerU 解析失败：{stderr.decode(errors='replace')[:500]}")

        # Find the largest .md file in the output tree
        md_files = list(out.rglob("*.md"))
        if not md_files:
            raise RuntimeError("MinerU 未生成 Markdown 文件")
        best = max(md_files, key=lambda p: p.stat().st_size)
        return best.read_text(encoding="utf-8", errors="replace")


# --------------------------------------------------------------------------- #
# python-docx
# --------------------------------------------------------------------------- #

def _docx(file_bytes: bytes) -> str:
    try:
        from docx import Document  # type: ignore
        import io

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return "\n".join(paragraphs)
    except ImportError:
        raise ImportError("请安装 python-docx：pip install python-docx")


# --------------------------------------------------------------------------- #
# python-pptx
# --------------------------------------------------------------------------- #

def _pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
        import io

        prs = Presentation(io.BytesIO(file_bytes))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
        return "\n".join(texts)
    except ImportError:
        raise ImportError("请安装 python-pptx：pip install python-pptx")


# --------------------------------------------------------------------------- #
# pypdf
# --------------------------------------------------------------------------- #

def _pdf(file_bytes: bytes) -> str:
    try:
        import pypdf  # type: ignore
        import io

        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except ImportError:
        raise ImportError("请安装 pypdf：pip install pypdf")


# --------------------------------------------------------------------------- #
# Vision model (image files)
# --------------------------------------------------------------------------- #

async def _image_vision(file_bytes: bytes, suffix: str) -> str:
    """Use the configured LLM vision API to OCR an image."""
    import instructor  # type: ignore
    from pydantic import BaseModel

    mime_map = {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".webp": "image/webp",
    }
    mime = mime_map.get(suffix, "image/jpeg")
    b64 = base64.b64encode(file_bytes).decode()

    class OCRResult(BaseModel):
        text: str

    client = _get_instructor_client()
    result: OCRResult = await asyncio.get_event_loop().run_in_executor(
        None,
        lambda: client.chat.completions.create(
            model=_get_model(),
            response_model=OCRResult,
            max_retries=2,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime};base64,{b64}"},
                        },
                        {
                            "type": "text",
                            "text": "请将图片中的所有文字原文提取出来，保留段落结构，不要添加任何说明。",
                        },
                    ],
                }
            ],
        ),
    )
    return result.text


# --------------------------------------------------------------------------- #
# Shared helpers (duplicated from pipeline to avoid circular import)
# --------------------------------------------------------------------------- #

def _get_instructor_client():  # type: ignore
    import instructor  # type: ignore

    provider = os.environ.get("LLM_PROVIDER", "deepseek")

    if provider == "deepseek":
        from openai import OpenAI  # type: ignore
        raw = OpenAI(
            api_key=os.environ["DEEPSEEK_API_KEY"],
            base_url="https://api.deepseek.com",
        )
        return instructor.from_openai(raw)

    if provider == "qwen":
        from openai import OpenAI  # type: ignore
        raw = OpenAI(
            api_key=os.environ["DASHSCOPE_API_KEY"],
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
        )
        return instructor.from_openai(raw)

    if provider == "openai":
        from openai import OpenAI  # type: ignore
        raw = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
        return instructor.from_openai(raw)

    raise ValueError(f"未知 LLM_PROVIDER: {provider}")


def _get_model() -> str:
    return os.environ.get("LLM_MODEL", "deepseek-chat")
